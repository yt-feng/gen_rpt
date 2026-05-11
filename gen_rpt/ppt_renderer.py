from __future__ import annotations

from pathlib import Path
from typing import Dict, Iterable, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .theme import load_theme

THEME = load_theme()
PALETTE = THEME["palette"]
BRAND_NAME = THEME.get("brand_name", "BlueOcean")

NAVY = PALETTE.get("navy_dark", "#051C2C")
ACCENT = PALETTE.get("accent", "#003087")
BRIGHT_BLUE = PALETTE.get("bright_blue", "#3273F6")
INK = PALETTE.get("ink", "#333333")
MUTED = PALETTE.get("subtle", "#64696E")
LINE = PALETTE.get("line", "#C8C8C8")
PANEL = PALETTE.get("panel", "#F5F7FA")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
LM = Inches(0.62)
RM = Inches(0.62)
TOP = Inches(0.38)


def render_pptx(report: Dict, assets: Dict[str, str], output_file: Path, topic: str, language: str = "en") -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    root = output_file.parent

    title = report.get("report_title", topic)
    subtitle = report.get("report_subtitle", "")

    _slide_cover(prs, title, subtitle, topic, assets, root)
    _slide_toc(prs, report)
    _slide_highlights(prs, report)

    for idx, section in enumerate(report.get("sections", [])[:10], start=1):
        _slide_section(prs, section, assets, root, idx)

    chart_paths = [v for k, v in _sorted_asset_items(assets) if k.startswith("chart-")]
    for idx, chart_path in enumerate(chart_paths[:6], start=1):
        _slide_chart(prs, f"Exhibit {idx}: evidence clarifies the strategic trade-off", chart_path, root, assets)

    _slide_closing(prs, report, topic, assets, root)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    return output_file


def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb("#FFFFFF")
    return slide


def _add_brand(slide, page_title: str = "") -> None:
    tx = slide.shapes.add_textbox(SLIDE_W - Inches(1.65), Inches(0.18), Inches(1.22), Inches(0.24))
    p = tx.text_frame.paragraphs[0]
    p.text = BRAND_NAME
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = _rgb(ACCENT)
    p.alignment = PP_ALIGN.RIGHT

    foot = slide.shapes.add_textbox(LM, SLIDE_H - Inches(0.28), SLIDE_W - LM - RM, Inches(0.16))
    p2 = foot.text_frame.paragraphs[0]
    p2.text = _fit(page_title, 120)
    p2.font.size = Pt(6)
    p2.font.color.rgb = _rgb("#A0A6AD")


def _title(slide, title: str, subtitle: str = "") -> None:
    box = slide.shapes.add_textbox(LM, TOP, Inches(10.6), Inches(0.92))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _fit(title, 118)
    p.font.size = Pt(20)
    p.font.color.rgb = _rgb(INK)
    p.font.bold = False
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = _fit(subtitle, 170)
        p2.font.size = Pt(8.3)
        p2.font.color.rgb = _rgb(MUTED)
    _line(slide, LM, Inches(1.22), SLIDE_W - LM - RM)


def _slide_cover(prs, title: str, subtitle: str, topic: str, assets: Dict[str, str], root: Path) -> None:
    slide = _blank(prs)
    bg_path = _asset_path(root, assets.get("cover-background", ""))
    if bg_path:
        _safe_picture(slide, bg_path, 0, 0, SLIDE_W, SLIDE_H)
        shade = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        shade.fill.solid(); shade.fill.fore_color.rgb = _rgb(NAVY); shade.fill.transparency = 18
        shade.line.fill.background()
    else:
        bg = slide.background.fill
        bg.solid(); bg.fore_color.rgb = _rgb(NAVY)

    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(0.75), Inches(5.35), Inches(3.75))
    panel.fill.solid(); panel.fill.fore_color.rgb = _rgb("#FFFFFF")
    panel.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(0.75), Inches(5.35), Inches(0.08))
    accent.fill.solid(); accent.fill.fore_color.rgb = _rgb(BRIGHT_BLUE); accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(1.05), Inches(1.05), Inches(4.65), Inches(3.0))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = BRAND_NAME.upper()
    p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = _rgb(ACCENT)
    p2 = tf.add_paragraph(); p2.text = _fit(title, 105); p2.font.size = Pt(23); p2.font.color.rgb = _rgb(INK); p2.space_before = Pt(18)
    p3 = tf.add_paragraph(); p3.text = _fit(subtitle or topic, 150); p3.font.size = Pt(9.5); p3.font.color.rgb = _rgb(MUTED); p3.space_before = Pt(14)


def _slide_toc(prs, report: Dict) -> None:
    slide = _blank(prs); _add_brand(slide, "Contents"); _title(slide, "The discussion follows the management question, evidence, and implications")
    sections = report.get("sections", [])[:10]
    if not sections:
        sections = [{"title": "Executive priorities and implications"}, {"title": "Evidence base and management agenda"}]
    y = Inches(1.55)
    for idx, sec in enumerate(sections, start=1):
        pnum = slide.shapes.add_textbox(LM, y, Inches(0.55), Inches(0.3)).text_frame.paragraphs[0]
        pnum.text = f"{idx:02d}"; pnum.font.size = Pt(9.5); pnum.font.bold = True; pnum.font.color.rgb = _rgb(ACCENT)
        p = slide.shapes.add_textbox(Inches(1.25), y, Inches(10.5), Inches(0.42)).text_frame.paragraphs[0]
        p.text = _fit(sec.get("title", "Section"), 110); p.font.size = Pt(12.6); p.font.color.rgb = _rgb(INK)
        y += Inches(0.50)


def _slide_highlights(prs, report: Dict) -> None:
    slide = _blank(prs); _add_brand(slide, "Key highlights"); _title(slide, "The analysis narrows the agenda to a focused set of management priorities")
    items = report.get("executive_summary", [])[:8]
    if not items:
        items = ["Evidence should be converted into a focused management agenda.", "Priorities should be sequenced by urgency, proof quality, and execution risk."]
    x0, y0 = LM, Inches(1.52)
    w, h = Inches(5.9), Inches(0.95)
    for idx, item in enumerate(items[:8]):
        x = x0 + Inches(6.12) * (idx % 2)
        y = y0 + Inches(1.11) * (idx // 2)
        _card(slide, x, y, w, h, _fit(_clean_summary_item(item), 135), idx + 1)


def _slide_section(prs, section: Dict, assets: Dict[str, str], root: Path, section_idx: int) -> None:
    slide = _blank(prs); _add_brand(slide, section.get("title", "Section")); _title(slide, section.get("title", "Section"), section.get("lead", ""))
    left = slide.shapes.add_textbox(LM, Inches(1.48), Inches(5.35), Inches(4.72))
    tf = left.text_frame; tf.word_wrap = True; tf.margin_top = 0; tf.margin_bottom = 0
    paragraphs = section.get("paragraphs", [])[:3] or [section.get("lead", "Evidence should be translated into management implications.")]
    for idx, paragraph in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = _fit(paragraph, 260)
        p.font.size = Pt(8.2)
        p.font.color.rgb = _rgb(INK)
        p.space_after = Pt(6)
    takeaways = section.get("key_takeaways", [])[:2]
    if takeaways:
        p = tf.add_paragraph(); p.text = "Key implications"; p.font.size = Pt(8.2); p.font.bold = True; p.font.color.rgb = _rgb(ACCENT); p.space_before = Pt(6)
        for item in takeaways:
            p2 = tf.add_paragraph(); p2.text = f"- {_fit(item, 96)}"; p2.font.size = Pt(7.3); p2.font.color.rgb = _rgb(INK)

    visual_path = _resolve_section_visual(section, section_idx, assets, root)
    if visual_path:
        _safe_picture(slide, visual_path, Inches(6.45), Inches(1.50), Inches(6.18), Inches(4.58))
    else:
        _visual_placeholder(slide, Inches(6.55), Inches(1.62), Inches(5.8), Inches(4.1), section.get("title", "Evidence visual"))


def _slide_chart(prs, title: str, chart_path: str, root: Path, assets: Dict[str, str]) -> None:
    slide = _blank(prs); _add_brand(slide, title); _title(slide, title)
    full = _asset_path(root, chart_path) or _asset_path(root, assets.get("cover-background", ""))
    if full:
        _safe_picture(slide, full, Inches(0.95), Inches(1.30), Inches(11.35), Inches(5.30))
    else:
        _visual_placeholder(slide, Inches(0.95), Inches(1.30), Inches(11.35), Inches(5.30), "Evidence exhibit")


def _slide_closing(prs, report: Dict, topic: str, assets: Dict[str, str], root: Path) -> None:
    slide = _blank(prs); _add_brand(slide, "Closing")
    cover = _asset_path(root, assets.get("cover-background", ""))
    if cover:
        _safe_picture(slide, cover, Inches(7.2), Inches(1.15), Inches(5.55), Inches(5.25))
    _title(slide, "The next step is to convert the research into a short list of decisions")
    tx = slide.shapes.add_textbox(LM, Inches(1.6), Inches(6.25), Inches(3.2))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _fit(report.get("report_subtitle", topic), 260)
    p.font.size = Pt(15); p.font.color.rgb = _rgb(ACCENT)


def _card(slide, x, y, w, h, text: str, idx: int) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = _rgb("#FFFFFF")
    box.line.color.rgb = _rgb(LINE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(ACCENT)
    bar.line.fill.background()
    num = slide.shapes.add_textbox(x + Inches(0.16), y + Inches(0.10), Inches(0.48), Inches(0.22))
    p = num.text_frame.paragraphs[0]; p.text = f"{idx:02d}"; p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = _rgb(ACCENT)
    body = slide.shapes.add_textbox(x + Inches(0.70), y + Inches(0.12), w - Inches(0.88), h - Inches(0.18))
    body.text_frame.word_wrap = True
    p2 = body.text_frame.paragraphs[0]; p2.text = text; p2.font.size = Pt(7.7); p2.font.color.rgb = _rgb(INK)


def _line(slide, x, y, w) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.01))
    line.fill.solid(); line.fill.fore_color.rgb = _rgb(LINE)
    line.line.color.rgb = _rgb(LINE)


def _resolve_section_visual(section: Dict, section_idx: int, assets: Dict[str, str], root: Path) -> Path | None:
    candidates = [
        f"image-{section_idx}",
        str(section.get("visual_hint", "")),
        f"chart-{section_idx}",
        "cover-background",
    ]
    for key in candidates:
        if not key:
            continue
        path = _asset_path(root, assets.get(key, ""))
        if path:
            return path
    return None


def _asset_path(root: Path, relative: str) -> Path | None:
    if not relative:
        return None
    full = root / relative
    if full.exists() and full.is_file() and full.stat().st_size > 0:
        return full
    return None


def _safe_picture(slide, path: Path, x, y, width, height) -> bool:
    try:
        slide.shapes.add_picture(str(path), x, y, width=width, height=height)
        return True
    except Exception:
        _visual_placeholder(slide, x, y, width, height, "Visual unavailable")
        return False


def _visual_placeholder(slide, x, y, w, h, label: str) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = _rgb(PANEL); box.line.color.rgb = _rgb(LINE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    accent.fill.solid(); accent.fill.fore_color.rgb = _rgb(ACCENT); accent.line.fill.background()
    tx = slide.shapes.add_textbox(x + Inches(0.28), y + Inches(0.28), w - Inches(0.55), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = _fit(label, 120)
    p.font.size = Pt(13)
    p.font.color.rgb = _rgb(ACCENT)


def _sorted_asset_items(assets: Dict[str, str]) -> Iterable[Tuple[str, str]]:
    def key(item: Tuple[str, str]):
        name = item[0]
        if "-" in name:
            try:
                return (name.split("-", 1)[0], int(name.split("-", 1)[1]))
            except Exception:
                pass
        return (name, 0)
    return sorted(assets.items(), key=key)


def _clean_summary_item(item: str) -> str:
    item = str(item or "").strip()
    if "：" in item:
        head, rest = item.split("：", 1)
        if rest.strip().startswith(head.strip()):
            item = head + "：" + rest.strip()[len(head.strip()):].lstrip("：: ，,。")
    return item


def _rgb(hex_color: str) -> RGBColor:
    color = hex_color.strip().lstrip("#")
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def _fit(text: str, limit: int) -> str:
    text = " ".join(str(text or "").split())
    if len(text) <= limit:
        return text
    return text[:limit].rstrip()
